"""Day 7. PowerPoint export for BRIEF.

Five slides:
  1. Cover (company + period + headline verdict).
  2. Tone curve PNG full bleed.
  3. Forward guidance table.
  4. Top Q&A exchanges.
  5. Risk flags + bull/bear case.

python-pptx is optional; is_available() returns False if not installed.
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

from analysis import CallAnalysis
from brief_chart import render_tone_curve_png

NAVY    = (0x0A, 0x16, 0x28)
CREAM   = (0xF4, 0xF1, 0xEA)
ORANGE  = (0xFF, 0x6B, 0x00)
CYAN    = (0x06, 0xB6, 0xD4)
GREY    = (0x6B, 0x72, 0x80)


def is_available() -> bool:
    return _HAS_PPTX


def write_pptx(call: CallAnalysis, out_path: Path) -> Path:
    if not _HAS_PPTX:
        raise RuntimeError("python-pptx not installed.")
    md = call.transcript.metadata
    h = call.headline
    es = call.exec_summary

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY, prs)
    _add_text(s, "BRIEF", Inches(0.6), Inches(0.5), Inches(4), Inches(0.6),
              size=24, bold=True, color=ORANGE)
    _add_text(s, "Day 07 . Earnings Call Summariser", Inches(0.6), Inches(1.0),
              Inches(8), Inches(0.4), size=12, color=GREY)
    _add_text(s, f"{md.company} ({md.ticker})", Inches(0.6), Inches(1.7),
              Inches(12), Inches(1.0), size=36, bold=True, color=CREAM)
    _add_text(s, f"{md.fiscal_period}  |  {md.call_date}", Inches(0.6),
              Inches(2.7), Inches(12), Inches(0.5), size=14, color=GREY)
    _add_text(s,
              f"Tone: {h.overall_tone.upper()}   "
              f"Confidence: {h.confidence_score:.2f}   "
              f"Hedge: {h.hedge_count}   "
              f"Certainty: {h.certainty_count}   "
              f"Deflection: {h.deflection_count}",
              Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
              size=14, color=CREAM)
    if es and es.headline:
        _add_text(s, es.headline, Inches(0.6), Inches(4.4), Inches(12), Inches(2),
                  size=18, color=CREAM)

    # Slide 2: Tone curve
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY, prs)
    _add_text(s, "Management tone curve", Inches(0.6), Inches(0.4),
              Inches(12), Inches(0.5), size=20, bold=True, color=CREAM)
    png = render_tone_curve_png(call.tone_curve)
    s.shapes.add_picture(BytesIO(png), Inches(0.6), Inches(1.2),
                         width=Inches(12), height=Inches(5.6))

    # Slide 3: Forward guidance
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY, prs)
    _add_text(s, "Forward guidance", Inches(0.6), Inches(0.4), Inches(12),
              Inches(0.5), size=20, bold=True, color=CREAM)
    y = Inches(1.2)
    _add_text(s, "Metric", Inches(0.6), y, Inches(2), Inches(0.4),
              size=11, bold=True, color=ORANGE)
    _add_text(s, "Period", Inches(2.6), y, Inches(2), Inches(0.4),
              size=11, bold=True, color=ORANGE)
    _add_text(s, "Direction", Inches(4.6), y, Inches(2), Inches(0.4),
              size=11, bold=True, color=ORANGE)
    _add_text(s, "Range", Inches(6.6), y, Inches(2), Inches(0.4),
              size=11, bold=True, color=ORANGE)
    _add_text(s, "Quote", Inches(8.6), y, Inches(4), Inches(0.4),
              size=11, bold=True, color=ORANGE)
    yi = y + Inches(0.4)
    for g in call.guidance[:10]:
        rng = ""
        if g.range_low is not None and g.range_high is not None:
            rng = f"{g.range_low:g} - {g.range_high:g} {g.unit}"
        elif g.range_low is not None:
            rng = f"{g.range_low:g} {g.unit}"
        _add_text(s, g.metric[:18], Inches(0.6), yi, Inches(2), Inches(0.4), size=10, color=CREAM)
        _add_text(s, g.period[:14], Inches(2.6), yi, Inches(2), Inches(0.4), size=10, color=CREAM)
        _add_text(s, g.direction[:10], Inches(4.6), yi, Inches(2), Inches(0.4), size=10, color=CREAM)
        _add_text(s, rng[:18], Inches(6.6), yi, Inches(2), Inches(0.4), size=10, color=CREAM)
        _add_text(s, g.quote[:80], Inches(8.6), yi, Inches(4), Inches(0.4), size=9, color=CREAM)
        yi += Inches(0.45)

    # Slide 4: Q&A
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY, prs)
    _add_text(s, "Analyst Q & A (top 6)", Inches(0.6), Inches(0.4),
              Inches(12), Inches(0.5), size=20, bold=True, color=CREAM)
    y = Inches(1.2)
    for qx in call.qa[:6]:
        _add_text(s, f"{qx.analyst_name} ({qx.analyst_firm})  -  {qx.tension} / {qx.management_clarity}",
                  Inches(0.6), y, Inches(12), Inches(0.4),
                  size=11, bold=True, color=ORANGE)
        y += Inches(0.35)
        _add_text(s, "Q: " + qx.question_summary[:180],
                  Inches(0.6), y, Inches(12), Inches(0.5), size=10, color=CREAM)
        y += Inches(0.4)
        _add_text(s, "A: " + qx.answer_summary[:240],
                  Inches(0.6), y, Inches(12), Inches(0.6),
                  size=10, color=(0x94, 0xA3, 0xB8))
        y += Inches(0.55)

    # Slide 5: Risks + bull/bear
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY, prs)
    _add_text(s, "Risks and verdict", Inches(0.6), Inches(0.4),
              Inches(12), Inches(0.5), size=20, bold=True, color=CREAM)
    y = Inches(1.2)
    if call.risks:
        _add_text(s, "Risk flags:", Inches(0.6), y, Inches(8), Inches(0.4),
                  size=12, bold=True, color=ORANGE)
        y += Inches(0.4)
        for r in call.risks[:6]:
            _add_text(s, f"[{r.severity.upper()}] {r.category}: {r.quote[:150]}",
                      Inches(0.6), y, Inches(12), Inches(0.5), size=10, color=CREAM)
            y += Inches(0.4)
    if es:
        if es.bull_case:
            _add_text(s, "Bull case:", Inches(0.6), y + Inches(0.2),
                      Inches(8), Inches(0.4), size=12, bold=True, color=CYAN)
            _add_text(s, es.bull_case, Inches(0.6), y + Inches(0.6),
                      Inches(12), Inches(0.8), size=10, color=CREAM)
            y += Inches(1.2)
        if es.bear_case:
            _add_text(s, "Bear case:", Inches(0.6), y + Inches(0.2),
                      Inches(8), Inches(0.4), size=12, bold=True, color=ORANGE)
            _add_text(s, es.bear_case, Inches(0.6), y + Inches(0.6),
                      Inches(12), Inches(0.8), size=10, color=CREAM)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def _bg(slide, rgb, prs):
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*rgb)
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_text(slide, text, left, top, width, height, *,
              size=12, bold=False, color=CREAM):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.color.rgb = RGBColor(*color)
    return box
